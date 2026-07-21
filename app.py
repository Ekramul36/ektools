from flask import Flask, render_template, request, send_file, send_from_directory, redirect, Response, jsonify, url_for, session
from pypdf import PdfReader, PdfWriter
import fitz
import os
import datetime
import uuid
import time
import shutil
from app_pdf_editor_route import pdf_editor_bp

app = Flask(__name__)
app.register_blueprint(pdf_editor_bp)
app.secret_key = os.environ.get("SECRET_KEY", "ekzapp-dev-secret-change-in-production")

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "outputs"
FAVICON_OUTPUT_DIR = os.path.join("static", "generated", "favicons")

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
os.makedirs(FAVICON_OUTPUT_DIR, exist_ok=True)


def unique_name(prefix, ext):
    """Generate a collision-safe output filename, e.g. Merged_PDF_a1b2c3d4.pdf
    This prevents one user's file from overwriting another user's file
    when multiple people use the same tool at the same time."""
    return f"{prefix}_{uuid.uuid4().hex[:8]}.{ext}"


FAVICON_MAX_AGE_HOURS = 24  # delete generated favicon folders older than this

def cleanup_old_favicons():
    """Delete favicon batch folders older than FAVICON_MAX_AGE_HOURS."""
    try:
        if not os.path.isdir(FAVICON_OUTPUT_DIR):
            return
        now = time.time()
        max_age_seconds = FAVICON_MAX_AGE_HOURS * 3600
        for folder_name in os.listdir(FAVICON_OUTPUT_DIR):
            folder_path = os.path.join(FAVICON_OUTPUT_DIR, folder_name)
            if not os.path.isdir(folder_path):
                continue
            folder_age = now - os.path.getmtime(folder_path)
            if folder_age > max_age_seconds:
                shutil.rmtree(folder_path, ignore_errors=True)
    except Exception:
        pass


# ──────────────────────────────────────────
# SEO — robots.txt & sitemap.xml & ads.txt
# ──────────────────────────────────────────

@app.route('/robots.txt')
def robots():
    content = "User-agent: *\nAllow: /\n\nSitemap: https://ekzapp.com/sitemap.xml"
    return Response(content, mimetype='text/plain')


@app.route('/ads.txt')
def ads_txt():
    content = "google.com, pub-3671123199046758, DIRECT, f08c47fec0942fa0"
    return Response(content, mimetype='text/plain')


# Site launch date — the fallback lastmod for any page not explicitly
# overridden below. Only bump LASTMOD_OVERRIDES for a page when you
# actually redesign/rewrite its content, not on every deploy.
LAUNCH_DATE = '2026-06-30'

LASTMOD_OVERRIDES = {
    '/': '2026-07-04',                          # homepage internal-link fix
    '/about': '2026-07-08',                     # E-E-A-T rewrite
    '/merge-pdf': '2026-07-06',                 # Tier 1 redesign
    '/compress-pdf': '2026-07-06',
    '/split-pdf': '2026-07-06',
    '/compress-image': '2026-07-06',
    '/remove-background': '2026-07-07',
    '/pdf-to-ppt': '2026-07-07',
    '/sign-pdf': '2026-07-08',                  # Tier 2 redesign
    '/crop-pdf': '2026-07-08',
    '/extract-images': '2026-07-08',
    '/passport-photo': '2026-07-08',
    '/meme-generator': '2026-07-10',            # content enrichment pass
    '/resume-builder': '2026-07-10',
    '/signature-generator': '2026-07-10',
    '/case-converter': '2026-07-10',
    '/blog/merge-pdf-files-online': '2026-07-01',
    '/blog/reduce-image-size-to-kb': '2026-07-01',
    '/blog/compress-pdf-to-100kb-for-exam-form': '2026-07-03',
    '/blog/amount-in-words-for-cheque-gst-invoice': '2026-07-03',
    '/blog/sign-pdf-online-free': '2026-07-03',
    '/blog/how-to-make-a-meme-online-free': '2026-07-06',
    '/blog/exam-photo-signature-size-guide-2026': '2026-07-08',
    '/blog/split-pdf-crop-pdf-guide': '2026-07-12',
    '/blog/unlock-pdf-remove-password-online-free': '2026-07-15',
    '/blog/itr-filing-2026-last-date-old-vs-new-tax-regime': '2026-07-17',
    '/blog/home-loan-emi-calculation-explained': '2026-07-19',
    '/cft-calculator': '2026-07-19',
    '/sand-calculator': '2026-07-19',
    '/blog/cement-sand-ratio-material-estimation-guide': '2026-07-20',
    '/land-area-converter': '2026-07-22',
    '/brick-calculator': '2026-07-22',
    '/in-hand-salary-calculator': '2026-07-22',
    '/hra-exemption-calculator': '2026-07-22',
    '/timestamp-converter': '2026-07-22',
}


@app.route('/sitemap.xml')
def sitemap():
    pages = [
        ('/', '1.0', 'daily'),
        ('/about', '0.8', 'monthly'),
        ('/contact', '0.8', 'monthly'),
        ('/privacy', '0.5', 'monthly'),
        ('/terms', '0.5', 'monthly'),
        ('/unlock-pdf', '0.9', 'weekly'),
        ('/merge-pdf', '0.9', 'weekly'),
        ('/split-pdf', '0.9', 'weekly'),
        ('/compress-pdf', '0.9', 'weekly'),
        ('/rotate-pdf', '0.9', 'weekly'),
        ('/delete-pages', '0.9', 'weekly'),
        ('/watermark-pdf', '0.9', 'weekly'),
        ('/protect-pdf', '0.9', 'weekly'),
        ('/pdf-to-jpg', '0.9', 'weekly'),
        ('/pdf-to-word', '0.9', 'weekly'),
        ('/word-to-pdf', '0.9', 'weekly'),
        ('/excel-to-pdf', '0.9', 'weekly'),
        ('/pdf-to-excel', '0.9', 'weekly'),
        ('/ocr', '0.9', 'weekly'),
        ('/pdf-editor', '0.9', 'weekly'),
        ('/crop-pdf', '0.9', 'weekly'),
        ('/extract-images', '0.9', 'weekly'),
        ('/pdf-page-reorder', '0.9', 'weekly'),
        ('/repair-pdf', '0.9', 'weekly'),
        ('/sign-pdf', '0.9', 'weekly'),
        ('/pdf-to-ppt', '0.9', 'weekly'),
        ('/image-to-pdf', '0.9', 'weekly'),
        ('/compress-image', '0.9', 'weekly'),
        ('/resize-image', '0.9', 'weekly'),
        ('/convert-image', '0.9', 'weekly'),
        ('/image-size-reducer', '0.9', 'weekly'),
        ('/remove-background', '0.9', 'weekly'),
        ('/passport-photo', '0.9', 'weekly'),
        ('/image-protect', '0.9', 'weekly'),
        ('/age-calculator', '0.9', 'weekly'),
        ('/bmi-calculator', '0.9', 'weekly'),
        ('/gst-calculator', '0.9', 'weekly'),
        ('/emi-calculator', '0.9', 'weekly'),
        ('/percentage-calculator', '0.9', 'weekly'),
        ('/sip-calculator', '0.9', 'weekly'),
        ('/date-calculator', '0.9', 'weekly'),
        ('/calorie-calculator', '0.9', 'weekly'),
        ('/tip-calculator', '0.9', 'weekly'),
        ('/word-counter', '0.9', 'weekly'),
        ('/case-converter', '0.9', 'weekly'),
        ('/mb-converter', '0.9', 'weekly'),
        ('/lorem-ipsum', '0.9', 'weekly'),
        ('/qr-generator', '0.9', 'weekly'),
        ('/password-generator', '0.9', 'weekly'),
        ('/currency-converter', '0.9', 'weekly'),
        ('/unit-converter', '0.9', 'weekly'),
        ('/cft-calculator', '0.9', 'weekly'),
        ('/sand-calculator', '0.9', 'weekly'),
        ('/invoice-generator', '0.9', 'weekly'),
        ('/youtube-thumbnail', '0.9', 'weekly'),
        ('/income-tax-calculator', '0.9', 'weekly'),
        ('/json-formatter', '0.9', 'weekly'),
        ('/color-picker', '0.9', 'weekly'),
        ('/signature-generator', '0.9', 'weekly'),
        ('/meme-generator', '0.9', 'weekly'),
        ('/resume-builder', '0.9', 'weekly'),
        ('/add-page-numbers', '0.9', 'weekly'),
        ('/amount-in-words', '0.9', 'weekly'),
        ('/base64-url-encoder', '0.9', 'weekly'),
        ('/favicon-generator', '0.9', 'weekly'),
        ('/land-area-converter', '0.9', 'weekly'),
        ('/brick-calculator', '0.9', 'weekly'),
        ('/in-hand-salary-calculator', '0.9', 'weekly'),
        ('/hra-exemption-calculator', '0.9', 'weekly'),
        ('/timestamp-converter', '0.9', 'weekly'),
        ('/blog', '0.8', 'weekly'),
        ('/blog/merge-pdf-files-online', '0.7', 'monthly'),
        ('/blog/reduce-image-size-to-kb', '0.7', 'monthly'),
        ('/blog/compress-pdf-to-100kb-for-exam-form', '0.7', 'monthly'),
        ('/blog/amount-in-words-for-cheque-gst-invoice', '0.7', 'monthly'),
        ('/blog/sign-pdf-online-free', '0.7', 'monthly'),
        ('/blog/how-to-make-a-meme-online-free', '0.7', 'monthly'),
        ('/blog/exam-photo-signature-size-guide-2026', '0.7', 'monthly'),
        ('/blog/split-pdf-crop-pdf-guide', '0.7', 'monthly'),
        ('/blog/unlock-pdf-remove-password-online-free', '0.7', 'monthly'),
        ('/blog/itr-filing-2026-last-date-old-vs-new-tax-regime', '0.7', 'monthly'),
        ('/blog/home-loan-emi-calculation-explained', '0.7', 'monthly'),
        ('/blog/cement-sand-ratio-material-estimation-guide', '0.7', 'monthly'),
        ('/author/ekramul-hoque', '0.5', 'monthly'),
    ]
    xml = ['<?xml version="1.0" encoding="UTF-8"?>',
           '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">']
    for path, priority, freq in pages:
        lastmod = LASTMOD_OVERRIDES.get(path, LAUNCH_DATE)
        xml.append(f'  <url>\n    <loc>https://ekzapp.com{path}</loc>\n    <lastmod>{lastmod}</lastmod>\n    <changefreq>{freq}</changefreq>\n    <priority>{priority}</priority>\n  </url>')
    xml.append('</urlset>')
    return Response('\n'.join(xml), mimetype='application/xml')


@app.route("/google1b7c2f73edcb5802.html")
def google_verification():
    return send_from_directory("static", "google1b7c2f73edcb5802.html")


# ──────────────────────────────────────────
# HOME
# ──────────────────────────────────────────

@app.route("/")
def home():
    return render_template("index.html")


# ──────────────────────────────────────────
# STATIC PAGES
# ──────────────────────────────────────────

@app.route("/about")
def about():
    return render_template("about.html")

@app.route("/privacy")
def privacy():
    return render_template("privacy.html")

@app.route("/terms")
def terms():
    return render_template("terms.html")

@app.route("/contact")
def contact():
    return render_template("contact.html")


# ──────────────────────────────────────────
# PDF TOOLS
# ──────────────────────────────────────────

@app.route("/unlock-pdf", methods=["GET", "POST"])
def unlock_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            password = request.form.get("password", "")
            if not pdf or pdf.filename == "":
                return render_template("pdf.html", error="Please select a PDF file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Unlocked", "pdf"))
            pdf.save(input_path)
            reader = PdfReader(input_path)
            if reader.is_encrypted:
                result = reader.decrypt(password)
                if result == 0:
                    return render_template("pdf.html", error="Wrong PDF Password!")
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(output_path, "wb") as output_file:
                writer.write(output_file)
            return send_file(output_path, as_attachment=True, download_name="Unlocked_" + pdf.filename)
        except Exception as e:
            return render_template("pdf.html", error=f"Error: {str(e)}")
    return render_template("pdf.html")


@app.route("/merge-pdf", methods=["GET", "POST"])
def merge_pdf():
    if request.method == "POST":
        try:
            pdf_files = request.files.getlist("pdfs")
            if len(pdf_files) < 2:
                return render_template("merge_pdf.html", error="Please select at least 2 PDF files.")
            writer = PdfWriter()
            for pdf in pdf_files:
                reader = PdfReader(pdf)
                for page in reader.pages:
                    writer.add_page(page)
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Merged_PDF", "pdf"))
            with open(output_path, "wb") as output_file:
                writer.write(output_file)
            return send_file(output_path, as_attachment=True, download_name="Merged_PDF.pdf")
        except Exception as e:
            return render_template("merge_pdf.html", error=f"{str(e)}")
    return render_template("merge_pdf.html")


@app.route("/split-pdf", methods=["GET", "POST"])
def split_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            pages = request.form.get("pages", "").strip()
            if not pdf or pdf.filename == "":
                return render_template("split_pdf.html", error="Please select a PDF file.")
            reader = PdfReader(pdf)
            writer = PdfWriter()
            if "-" in pages:
                start, end = pages.split("-")
                for i in range(int(start) - 1, int(end)):
                    writer.add_page(reader.pages[i])
            else:
                writer.add_page(reader.pages[int(pages) - 1])
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Split_PDF", "pdf"))
            with open(output_path, "wb") as f:
                writer.write(f)
            return send_file(output_path, as_attachment=True, download_name="Split_PDF.pdf")
        except Exception as e:
            return render_template("split_pdf.html", error=f"{str(e)}")
    return render_template("split_pdf.html")


@app.route("/compress-pdf", methods=["GET", "POST"])
def compress_pdf():
    pdf = request.files.get("pdf")
    if not pdf or pdf.filename == "":
        return render_template("compress_pdf.html", error="Please select a PDF.")
    input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
    output_path = os.path.join(OUTPUT_FOLDER, unique_name("Compressed", "pdf"))
    pdf.save(input_path)
    try:
        doc = fitz.open(input_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()
        return send_file(output_path, as_attachment=True, download_name="Compressed_" + pdf.filename)
    except Exception as e:
        return render_template("compress_pdf.html", error=str(e))
    return render_template("compress_pdf.html")


@app.route("/rotate-pdf", methods=["GET", "POST"])
def rotate_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            rotation = int(request.form.get("rotation", 90))
            if not pdf or pdf.filename == "":
                return render_template("rotate_pdf.html", error="Please select a PDF file.")
            reader = PdfReader(pdf)
            writer = PdfWriter()
            for page in reader.pages:
                page.rotate(rotation)
                writer.add_page(page)
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Rotated_PDF", "pdf"))
            with open(output_path, "wb") as f:
                writer.write(f)
            return send_file(output_path, as_attachment=True, download_name="Rotated_PDF.pdf")
        except Exception as e:
            return render_template("rotate_pdf.html", error=f"Error: {str(e)}")
    return render_template("rotate_pdf.html")


@app.route("/delete-pages", methods=["GET", "POST"])
def delete_pages():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            pages_input = request.form.get("pages", "").strip()
            if not pdf or pdf.filename == "":
                return render_template("delete_pages.html", error="Please select a PDF file.")
            reader = PdfReader(pdf)
            pages_to_delete = set()
            for part in pages_input.split(","):
                part = part.strip()
                if "-" in part:
                    s, e = part.split("-")
                    for p in range(int(s), int(e)+1):
                        pages_to_delete.add(p)
                else:
                    pages_to_delete.add(int(part))
            writer = PdfWriter()
            for i, page in enumerate(reader.pages):
                if (i+1) not in pages_to_delete:
                    writer.add_page(page)
            if len(writer.pages) == 0:
                return render_template("delete_pages.html", error="Cannot delete all pages!")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Deleted_Pages", "pdf"))
            with open(output_path, "wb") as f:
                writer.write(f)
            return send_file(output_path, as_attachment=True, download_name="Deleted_Pages.pdf")
        except Exception as e:
            return render_template("delete_pages.html", error=f"Error: {str(e)}")
    return render_template("delete_pages.html")


@app.route("/watermark-pdf", methods=["GET", "POST"])
def watermark_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            watermark = request.form.get("watermark", "CONFIDENTIAL")
            font_size = int(request.form.get("font_size", 50))
            if not pdf or pdf.filename == "":
                return render_template("watermark_pdf.html", error="Please select a PDF.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Watermarked", "pdf"))
            pdf.save(input_path)
            doc = fitz.open(input_path)
            for page in doc:
                page.insert_text(
                    (page.rect.width/2 - len(watermark)*font_size/4, page.rect.height/2),
                    watermark, fontsize=font_size, color=(0.5, 0.5, 0.5), rotate=45)
            doc.save(output_path)
            doc.close()
            return send_file(output_path, as_attachment=True, download_name="Watermarked.pdf")
        except Exception as e:
            return render_template("watermark_pdf.html", error=str(e))
    return render_template("watermark_pdf.html")


@app.route("/protect-pdf", methods=["GET", "POST"])
def protect_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            password = request.form.get("password", "")
            confirm = request.form.get("confirm_password", "")
            if not pdf or pdf.filename == "":
                return render_template("protect_pdf.html", error="Please select a PDF.")
            if password != confirm:
                return render_template("protect_pdf.html", error="Passwords do not match!")
            if not password:
                return render_template("protect_pdf.html", error="Please enter a password.")
            reader = PdfReader(pdf)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(password)
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Protected", "pdf"))
            with open(output_path, "wb") as f:
                writer.write(f)
            return send_file(output_path, as_attachment=True, download_name="Protected.pdf")
        except Exception as e:
            return render_template("protect_pdf.html", error=str(e))
    return render_template("protect_pdf.html")


@app.route("/pdf-to-jpg", methods=["GET", "POST"])
def pdf_to_jpg():
    import zipfile
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            quality = int(request.form.get("quality", 200))
            if not pdf or pdf.filename == "":
                return render_template("pdf_to_jpg.html", error="Please select a PDF file.")
            batch_id = uuid.uuid4().hex[:8]
            input_path = os.path.join(UPLOAD_FOLDER, f"{batch_id}_{pdf.filename}")
            pdf.save(input_path)
            doc = fitz.open(input_path)
            if doc.page_count == 1:
                page = doc[0]
                mat = fitz.Matrix(quality/72, quality/72)
                pix = page.get_pixmap(matrix=mat)
                output_path = os.path.join(OUTPUT_FOLDER, f"{batch_id}_page_1.jpg")
                pix.save(output_path)
                doc.close()
                return send_file(output_path, as_attachment=True, download_name="page_1.jpg")
            else:
                zip_path = os.path.join(OUTPUT_FOLDER, f"{batch_id}_PDF_to_JPG.zip")
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for i, page in enumerate(doc):
                        mat = fitz.Matrix(quality/72, quality/72)
                        pix = page.get_pixmap(matrix=mat)
                        img_path = os.path.join(OUTPUT_FOLDER, f"{batch_id}_page_{i+1}.jpg")
                        pix.save(img_path)
                        zipf.write(img_path, f"page_{i+1}.jpg")
                doc.close()
                return send_file(zip_path, as_attachment=True, download_name="PDF_to_JPG.zip")
        except Exception as e:
            return render_template("pdf_to_jpg.html", error=f"Error: {str(e)}")
    return render_template("pdf_to_jpg.html")


@app.route("/pdf-to-word", methods=["GET", "POST"])
def pdf_to_word():
    if request.method == "POST":
        try:
            from pdf2docx import Converter
            pdf = request.files.get("pdf")
            if not pdf or pdf.filename == "":
                return render_template("pdf_to_word.html", error="Please select a PDF file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Converted", "docx"))
            pdf.save(input_path)
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            return send_file(output_path, as_attachment=True, download_name="Converted.docx")
        except Exception as e:
            return render_template("pdf_to_word.html", error=f"Error: {str(e)}")
    return render_template("pdf_to_word.html")


@app.route("/word-to-pdf", methods=["GET", "POST"])
def word_to_pdf():
    if request.method == "POST":
        try:
            from docx2pdf import convert
            doc = request.files.get("file")
            if not doc or doc.filename == "":
                return render_template("word_to_pdf.html", error="Please select a Word file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{doc.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Converted", "pdf"))
            doc.save(input_path)
            convert(input_path, output_path)
            return send_file(output_path, as_attachment=True, download_name="Converted.pdf")
        except Exception as e:
            return render_template("word_to_pdf.html", error=f"Error: {str(e)}")
    return render_template("word_to_pdf.html")


@app.route("/excel-to-pdf", methods=["GET", "POST"])
def excel_to_pdf():
    if request.method == "POST":
        try:
            import openpyxl
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors
            xl = request.files.get("file")
            if not xl or xl.filename == "":
                return render_template("excel_to_pdf.html", error="Please select an Excel file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{xl.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Converted", "pdf"))
            xl.save(input_path)
            wb = openpyxl.load_workbook(input_path, data_only=True)
            ws = wb.active
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else "" for cell in row])
            if not data:
                return render_template("excel_to_pdf.html", error="Excel file is empty.")
            doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
            col_count = len(data[0])
            col_width = (landscape(A4)[0] - 40) / col_count if col_count > 0 else 100
            table = Table(data, colWidths=[col_width] * col_count)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#2563EB')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('FONTSIZE', (0,0), (-1,-1), 8),
                ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#F9FAFB')]),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#E5E7EB')),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('PADDING', (0,0), (-1,-1), 4),
            ]))
            doc.build([table])
            return send_file(output_path, as_attachment=True, download_name="Converted.pdf")
        except Exception as e:
            return render_template("excel_to_pdf.html", error=f"Error: {str(e)}")
    return render_template("excel_to_pdf.html")


@app.route("/pdf-to-excel", methods=["GET", "POST"])
def pdf_to_excel():
    if request.method == "POST":
        try:
            import pdfplumber
            import openpyxl
            pdf = request.files.get("pdf")
            if not pdf or pdf.filename == "":
                return render_template("pdf_to_excel.html", error="Please select a PDF file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Converted", "xlsx"))
            pdf.save(input_path)
            wb = openpyxl.Workbook()
            first = True
            with pdfplumber.open(input_path) as pdf_doc:
                for i, page in enumerate(pdf_doc.pages):
                    tables = page.extract_tables()
                    if tables:
                        for t_idx, table in enumerate(tables):
                            if first:
                                ws = wb.active
                                ws.title = f"Page{i+1}"
                                first = False
                            else:
                                ws = wb.create_sheet(f"Page{i+1}_T{t_idx+1}")
                            for row in table:
                                ws.append([cell if cell else "" for cell in row])
                    else:
                        text = page.extract_text()
                        if text:
                            if first:
                                ws = wb.active
                                ws.title = f"Page{i+1}"
                                first = False
                            else:
                                ws = wb.create_sheet(f"Page{i+1}")
                            for line in text.split('\n'):
                                ws.append([line])
            if first:
                return render_template("pdf_to_excel.html", error="No content found in PDF.")
            wb.save(output_path)
            return send_file(output_path, as_attachment=True, download_name="Converted.xlsx")
        except Exception as e:
            return render_template("pdf_to_excel.html", error=f"Error: {str(e)}")
    return render_template("pdf_to_excel.html")


@app.route("/ocr", methods=["GET", "POST"])
def ocr_tool():
    if request.method == "POST":
        try:
            import pytesseract
            from PIL import Image as PILImage
            img_file = request.files.get("image")
            if not img_file or img_file.filename == "":
                return render_template("ocr.html", error="Please select an image.")
            img = PILImage.open(img_file)
            text = pytesseract.image_to_string(img, lang='eng')
            if not text.strip():
                return render_template("ocr.html", error="No text found in image.", result="")
            return render_template("ocr.html", result=text)
        except Exception as e:
            return render_template("ocr.html", error=f"Error: {str(e)}")
    return render_template("ocr.html")


@app.route("/pdf-editor")
def pdf_editor():
    return render_template("pdf_editor.html")


# ──────────────────────────────────────────
# NEW PDF TOOLS — Crop, Extract Images, Page Reorder, Repair, Sign, PDF to PPT
# ──────────────────────────────────────────

@app.route("/crop-pdf", methods=["GET", "POST"])
def crop_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            top = float(request.form.get("crop_top", 0) or 0)
            bottom = float(request.form.get("crop_bottom", 0) or 0)
            left = float(request.form.get("crop_left", 0) or 0)
            right = float(request.form.get("crop_right", 0) or 0)
            if not pdf or pdf.filename == "":
                return render_template("crop_pdf.html", error="Please select a PDF file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Cropped", "pdf"))
            pdf.save(input_path)
            mm_to_pt = 2.8346
            doc = fitz.open(input_path)
            for page in doc:
                r = page.rect
                new_rect = fitz.Rect(
                    r.x0 + left * mm_to_pt,
                    r.y0 + top * mm_to_pt,
                    r.x1 - right * mm_to_pt,
                    r.y1 - bottom * mm_to_pt
                )
                if new_rect.width > 10 and new_rect.height > 10:
                    page.set_cropbox(new_rect)
            doc.save(output_path)
            doc.close()
            return send_file(output_path, as_attachment=True, download_name="Cropped_" + pdf.filename)
        except Exception as e:
            return render_template("crop_pdf.html", error=f"Error: {str(e)}")
    return render_template("crop_pdf.html")


@app.route("/extract-images", methods=["GET", "POST"])
def extract_images():
    import zipfile
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            if not pdf or pdf.filename == "":
                return render_template("extract_images.html", error="Please select a PDF file.")
            batch_id = uuid.uuid4().hex[:8]
            input_path = os.path.join(UPLOAD_FOLDER, f"{batch_id}_{pdf.filename}")
            pdf.save(input_path)
            doc = fitz.open(input_path)
            count = 0
            zip_path = os.path.join(OUTPUT_FOLDER, f"{batch_id}_Extracted_Images.zip")
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for page_index in range(len(doc)):
                    page = doc[page_index]
                    images = page.get_images(full=True)
                    for img_index, img in enumerate(images):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        ext = base_image["ext"]
                        count += 1
                        img_filename = f"page{page_index+1}_image{img_index+1}.{ext}"
                        img_path = os.path.join(OUTPUT_FOLDER, f"{batch_id}_{img_filename}")
                        with open(img_path, "wb") as f:
                            f.write(image_bytes)
                        zipf.write(img_path, img_filename)
                        os.remove(img_path)
            doc.close()
            if count == 0:
                return render_template("extract_images.html", error="No images found in this PDF.")
            return send_file(zip_path, as_attachment=True, download_name="Extracted_Images.zip")
        except Exception as e:
            return render_template("extract_images.html", error=f"Error: {str(e)}")
    return render_template("extract_images.html")


@app.route("/pdf-page-reorder", methods=["GET"])
def pdf_page_reorder():
    return render_template("pdf_page_reorder.html")


@app.route("/pdf-page-reorder/upload", methods=["POST"])
def pdf_page_reorder_upload():
    try:
        pdf = request.files.get("pdf")
        if not pdf or pdf.filename == "":
            return jsonify({"error": "Please select a PDF file."}), 400
        token = uuid.uuid4().hex
        input_path = os.path.join(UPLOAD_FOLDER, f"reorder_{token}.pdf")
        pdf.save(input_path)
        reader = PdfReader(input_path)
        num_pages = len(reader.pages)
        return jsonify({"token": token, "num_pages": num_pages, "filename": pdf.filename})
    except Exception as e:
        return jsonify({"error": f"Error: {str(e)}"}), 500


@app.route("/pdf-page-reorder/apply", methods=["POST"])
def pdf_page_reorder_apply():
    try:
        token = request.form.get("token", "")
        order = request.form.get("order", "")
        input_path = os.path.join(UPLOAD_FOLDER, f"reorder_{token}.pdf")
        if not token or not os.path.exists(input_path):
            return render_template("pdf_page_reorder.html", error="Your session expired. Please upload the PDF again.")
        order_list = [int(x.strip()) - 1 for x in order.split(",") if x.strip() != ""]
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for idx in order_list:
            if 0 <= idx < len(reader.pages):
                writer.add_page(reader.pages[idx])
        if len(writer.pages) == 0:
            return render_template("pdf_page_reorder.html", error="Something went wrong with the page order. Please try again.")
        output_path = os.path.join(OUTPUT_FOLDER, unique_name("Reordered", "pdf"))
        with open(output_path, "wb") as f:
            writer.write(f)
        try:
            os.remove(input_path)
        except Exception:
            pass
        return send_file(output_path, as_attachment=True, download_name="Reordered.pdf")
    except Exception as e:
        return render_template("pdf_page_reorder.html", error=f"Error: {str(e)}")


@app.route("/repair-pdf", methods=["GET", "POST"])
def repair_pdf():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            if not pdf or pdf.filename == "":
                return render_template("repair_pdf.html", error="Please select a PDF file.")
            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Repaired", "pdf"))
            pdf.save(input_path)
            recovered_pages = 0
            try:
                doc = fitz.open(input_path)
                recovered_pages = len(doc)
                doc.save(output_path, garbage=4, deflate=True, clean=True)
                doc.close()
            except Exception:
                reader = PdfReader(input_path, strict=False)
                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)
                recovered_pages = len(writer.pages)
                with open(output_path, "wb") as f:
                    writer.write(f)
            if recovered_pages == 0:
                return render_template("repair_pdf.html", error="Could not recover any pages from this PDF. The file may be too severely damaged.")
            return send_file(output_path, as_attachment=True, download_name="Repaired_" + pdf.filename)
        except Exception as e:
            return render_template("repair_pdf.html", error=f"This PDF appears too damaged to repair automatically. Error: {str(e)}")
    return render_template("repair_pdf.html")


@app.route("/sign-pdf", methods=["GET", "POST"])
def sign_pdf():
    if request.method == "POST":
        try:
            import base64
            from PIL import Image as PILImage

            pdf = request.files.get("pdf")
            signature_data = request.form.get("signature_data", "")
            page_choice = request.form.get("page_choice", "last")
            position = request.form.get("position", "bottom-right")
            sig_width_pct = float(request.form.get("sig_width", 25) or 25)

            if not pdf or pdf.filename == "":
                return render_template("sign_pdf.html", error="Please select a PDF file.")
            if not signature_data or "," not in signature_data:
                return render_template("sign_pdf.html", error="Please draw or type your signature first.")

            sig_bytes = base64.b64decode(signature_data.split(",")[1])
            sig_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_sig.png")
            with open(sig_path, "wb") as f:
                f.write(sig_bytes)

            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Signed", "pdf"))
            pdf.save(input_path)

            doc = fitz.open(input_path)
            total = len(doc)

            if page_choice == "first":
                target_pages = [0]
            elif page_choice == "all":
                target_pages = list(range(total))
            else:
                target_pages = [total - 1]

            sig_pil = PILImage.open(sig_path)
            sig_ratio = sig_pil.height / sig_pil.width
            margin = 24

            for i in target_pages:
                page = doc[i]
                r = page.rect
                sig_w = r.width * (sig_width_pct / 100)
                sig_h = sig_w * sig_ratio

                if position == "bottom-left":
                    rect = fitz.Rect(margin, r.height - margin - sig_h, margin + sig_w, r.height - margin)
                elif position == "top-left":
                    rect = fitz.Rect(margin, margin, margin + sig_w, margin + sig_h)
                elif position == "top-right":
                    rect = fitz.Rect(r.width - margin - sig_w, margin, r.width - margin, margin + sig_h)
                elif position == "center":
                    rect = fitz.Rect(r.width/2 - sig_w/2, r.height/2 - sig_h/2, r.width/2 + sig_w/2, r.height/2 + sig_h/2)
                else:
                    rect = fitz.Rect(r.width - margin - sig_w, r.height - margin - sig_h, r.width - margin, r.height - margin)

                page.insert_image(rect, filename=sig_path)

            doc.save(output_path)
            doc.close()
            try:
                os.remove(sig_path)
            except Exception:
                pass
            return send_file(output_path, as_attachment=True, download_name="Signed_" + pdf.filename)
        except Exception as e:
            return render_template("sign_pdf.html", error=f"Error: {str(e)}")
    return render_template("sign_pdf.html")


@app.route("/pdf-to-ppt", methods=["GET", "POST"])
def pdf_to_ppt():
    if request.method == "POST":
        try:
            from pptx import Presentation
            from pptx.util import Emu

            pdf = request.files.get("pdf")
            if not pdf or pdf.filename == "":
                return render_template("pdf_to_ppt.html", error="Please select a PDF file.")

            batch_id = uuid.uuid4().hex[:8]
            input_path = os.path.join(UPLOAD_FOLDER, f"{batch_id}_{pdf.filename}")
            pdf.save(input_path)

            doc = fitz.open(input_path)
            if len(doc) == 0:
                return render_template("pdf_to_ppt.html", error="This PDF has no pages.")

            prs = Presentation()
            first_page = doc[0]
            emu_per_pt = 12700
            prs.slide_width = Emu(int(first_page.rect.width * emu_per_pt))
            prs.slide_height = Emu(int(first_page.rect.height * emu_per_pt))
            blank_layout = prs.slide_layouts[6]

            for i, page in enumerate(doc):
                mat = fitz.Matrix(2, 2)
                pix = page.get_pixmap(matrix=mat)
                img_path = os.path.join(UPLOAD_FOLDER, f"{batch_id}_slide_{i+1}.png")
                pix.save(img_path)
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                os.remove(img_path)

            doc.close()
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Converted", "pptx"))
            prs.save(output_path)
            return send_file(output_path, as_attachment=True, download_name="Converted.pptx")
        except Exception as e:
            return render_template("pdf_to_ppt.html", error=f"Error: {str(e)}")
    return render_template("pdf_to_ppt.html")


# ──────────────────────────────────────────
# IMAGE TOOLS
# ──────────────────────────────────────────

@app.route("/image-to-pdf", methods=["GET", "POST"])
def image_to_pdf():
    if request.method == "POST":
        try:
            from PIL import Image as PILImage
            images = request.files.getlist("images")
            if not images or images[0].filename == "":
                return render_template("image.html", error="Please select image files.")
            pdf_images = []
            for img_file in images:
                img = PILImage.open(img_file).convert("RGB")
                pdf_images.append(img)
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Images_to_PDF", "pdf"))
            if len(pdf_images) == 1:
                pdf_images[0].save(output_path, "PDF")
            else:
                pdf_images[0].save(output_path, "PDF", save_all=True, append_images=pdf_images[1:])
            return send_file(output_path, as_attachment=True, download_name="Images_to_PDF.pdf")
        except Exception as e:
            return render_template("image.html", error=f"Error: {str(e)}")
    return render_template("image.html")

@app.route("/image", methods=["GET", "POST"])
def image():
    return redirect("/image-to-pdf", code=301)


@app.route("/compress-image", methods=["GET", "POST"])
def compress_image():
    if request.method == "POST":
        try:
            from PIL import Image as PILImage
            import io as io_module
            import base64

            img_file = request.files.get("image")
            quality = int(request.form.get("quality", 80))
            if not img_file or img_file.filename == "":
                return render_template("image_compress.html", error="Please select an image.")

            # Measure original size in KB before touching the file
            img_file.seek(0, os.SEEK_END)
            original_size_kb = img_file.tell() / 1024
            img_file.seek(0)

            img = PILImage.open(img_file)
            if img.mode in ("RGBA", "LA", "P"):
                bg = PILImage.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                bg.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = bg
            elif img.mode != "RGB":
                img = img.convert("RGB")

            buf = io_module.BytesIO()
            img.save(buf, format="JPEG", quality=quality, optimize=True)
            buf.seek(0)
            final_bytes = buf.getvalue()
            final_size_kb = len(final_bytes) / 1024

            img_b64 = base64.b64encode(final_bytes).decode("utf-8")
            data_uri = f"data:image/jpeg;base64,{img_b64}"
            fname = img_file.filename.rsplit(".", 1)[0]

            reduction_pct = 0
            if original_size_kb > 0:
                reduction_pct = round((1 - (final_size_kb / original_size_kb)) * 100, 1)
                if reduction_pct < 0:
                    reduction_pct = 0

            return render_template(
                "image_compress.html",
                result_image=data_uri,
                original_size_kb=round(original_size_kb, 1),
                final_size_kb=round(final_size_kb, 1),
                reduction_pct=reduction_pct,
                quality=quality,
                download_filename=f"{fname}_compressed.jpg",
            )
        except Exception as e:
            return render_template("image_compress.html", error=str(e))
    return render_template("image_compress.html")

@app.route("/image/compress", methods=["GET", "POST"])
def image_compress_old():
    return redirect("/compress-image", code=301)


@app.route("/resize-image", methods=["GET", "POST"])
def resize_image():
    if request.method == "POST":
        try:
            from PIL import Image as PILImage
            img_file = request.files.get("image")
            width = request.form.get("width")
            height = request.form.get("height")
            keep_ratio = request.form.get("keep_ratio")
            if not img_file or img_file.filename == "":
                return render_template("image_resize.html", error="Please select an image.")
            img = PILImage.open(img_file)
            w = int(width) if width else None
            h = int(height) if height else None
            if keep_ratio and w and not h:
                h = int(img.height * (w / img.width))
            elif keep_ratio and h and not w:
                w = int(img.width * (h / img.height))
            if w and h:
                img = img.resize((w, h), PILImage.LANCZOS)
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Resized", "jpg"))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.save(output_path, "JPEG", quality=95)
            return send_file(output_path, as_attachment=True, download_name="Resized.jpg")
        except Exception as e:
            return render_template("image_resize.html", error=str(e))
    return render_template("image_resize.html")

@app.route("/image/resize", methods=["GET", "POST"])
def image_resize_old():
    return redirect("/resize-image", code=301)


@app.route("/convert-image", methods=["GET", "POST"])
def convert_image():
    if request.method == "POST":
        try:
            from PIL import Image as PILImage
            img_file = request.files.get("image")
            fmt = request.form.get("format", "jpg").upper()
            if not img_file or img_file.filename == "":
                return render_template("image_convert.html", error="Please select an image.")
            img = PILImage.open(img_file)
            if fmt == "JPG":
                fmt = "JPEG"
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
            ext = "jpg" if fmt == "JPEG" else fmt.lower()
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Converted", ext))
            img.save(output_path, fmt)
            return send_file(output_path, as_attachment=True, download_name=f"Converted.{ext}")
        except Exception as e:
            return render_template("image_convert.html", error=str(e))
    return render_template("image_convert.html")

@app.route("/image/convert", methods=["GET", "POST"])
def image_convert_old():
    return redirect("/convert-image", code=301)


@app.route("/image-size-reducer", methods=["GET", "POST"])
def image_size_reducer():
    if request.method == "POST":
        try:
            from PIL import Image as PILImage
            import io as io_module
            import base64

            img_file = request.files.get("image")
            target_kb = request.form.get("target_kb", type=int)

            if not img_file or img_file.filename == "":
                return render_template("image_size_reducer.html", error="Please select an image.")
            if not target_kb or target_kb < 5:
                return render_template("image_size_reducer.html", error="Target size must be at least 5 KB.")

            # Measure the original file size in KB
            img_file.seek(0, os.SEEK_END)
            original_size_kb = img_file.tell() / 1024
            img_file.seek(0)

            target_bytes = target_kb * 1024
            img = PILImage.open(img_file)

            # Flatten transparency / palette / CMYK to white-background RGB
            if img.mode in ("RGBA", "LA", "P"):
                bg = PILImage.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                bg.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = bg
            elif img.mode != "RGB":
                img = img.convert("RGB")

            ow, oh = img.size
            best_buffer = None

            # Step 1: binary-search JPEG quality at full resolution
            low, high = 5, 95
            while low <= high:
                mid = (low + high) // 2
                buf = io_module.BytesIO()
                img.save(buf, format="JPEG", quality=mid, optimize=True)
                if buf.tell() <= target_bytes:
                    best_buffer = buf
                    low = mid + 1
                else:
                    high = mid - 1

            # Step 2: if still too big at lowest quality, shrink resolution
            if best_buffer is None:
                scale = 0.9
                buf = None
                while scale > 0.1:
                    resized = img.resize((max(1, int(ow * scale)), max(1, int(oh * scale))), PILImage.LANCZOS)
                    buf = io_module.BytesIO()
                    resized.save(buf, format="JPEG", quality=5, optimize=True)
                    if buf.tell() <= target_bytes:
                        best_buffer = buf
                        break
                    scale -= 0.1
                if best_buffer is None:
                    best_buffer = buf  # last attempt, smallest we could get

            best_buffer.seek(0)
            final_size_kb = len(best_buffer.getvalue()) / 1024
            img_b64 = base64.b64encode(best_buffer.getvalue()).decode("utf-8")
            data_uri = f"data:image/jpeg;base64,{img_b64}"
            fname = img_file.filename.rsplit(".", 1)[0]

            return render_template(
                "image_size_reducer.html",
                result_image=data_uri,
                original_size_kb=round(original_size_kb, 1),
                final_size_kb=round(final_size_kb, 1),
                target_kb=target_kb,
                download_filename=f"{fname}_compressed.jpg",
            )
        except Exception as e:
            return render_template("image_size_reducer.html", error=f"Error: {str(e)}")
    return render_template("image_size_reducer.html")


@app.route("/remove-background", methods=["GET", "POST"])
def remove_background():
    if request.method == "POST":
        try:
            import requests as req_lib
            img_file = request.files.get("image")
            if not img_file or img_file.filename == "":
                return render_template("remove_background.html", error="Please select an image.")
            response = req_lib.post(
                'https://api.remove.bg/v1.0/removebg',
                files={'image_file': img_file.read()},
                data={'size': 'auto'},
                headers={'X-Api-Key': os.environ.get('REMOVE_BG_API_KEY')},
            )
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("No_Background", "png"))
            with open(output_path, "wb") as f:
                f.write(response.content)
            return send_file(output_path, as_attachment=True, download_name="No_Background.png")
        except Exception as e:
            return render_template("remove_background.html", error=f"Error: {str(e)}")
    return render_template("remove_background.html")


@app.route("/passport-photo", methods=["GET", "POST"])
def passport_photo():
    if request.method == "POST":
        try:
            from PIL import Image as PILImage
            import base64
            import io as BytesIO_io

            img_file    = request.files.get("image")
            size_preset = request.form.get("size", "35x45")
            bg_color    = request.form.get("bg_color", "#ffffff")
            copies      = int(request.form.get("copies", 8))

            if not img_file or img_file.filename == "":
                return render_template("passport_photo.html", error="Please select an image.")

            sizes = {
                "35x45": (413, 531),
                "51x51": (600, 600),
                "35x35": (413, 413),
                "40x60": (472, 709),
            }
            w, h = sizes.get(size_preset, (413, 531))
            img = PILImage.open(img_file).convert("RGBA")
            img_ratio    = img.width / img.height
            target_ratio = w / h
            if img_ratio > target_ratio:
                new_h = h; new_w = int(h * img_ratio)
            else:
                new_w = w; new_h = int(w / img_ratio)
            img  = img.resize((new_w, new_h), PILImage.LANCZOS)
            left = (new_w - w) // 2; top = (new_h - h) // 2
            img  = img.crop((left, top, left + w, top + h))
            try:
                bg = PILImage.new("RGB", (w, h), bg_color)
            except Exception:
                bg = PILImage.new("RGB", (w, h), "white")
            if img.mode == "RGBA":
                bg.paste(img, mask=img.split()[3])
            else:
                bg.paste(img)

            dpi = 300
            a4_w = int(8.27 * dpi); a4_h = int(11.69 * dpi)
            sheet = PILImage.new("RGB", (a4_w, a4_h), "white")
            margin = int(0.3 * dpi); gap = int(0.1 * dpi)
            placed = 0; x_pos = margin; y_pos = margin
            while placed < copies:
                if x_pos + w > a4_w - margin:
                    x_pos = margin; y_pos += h + gap
                if y_pos + h > a4_h - margin:
                    break
                sheet.paste(bg, (x_pos, y_pos))
                x_pos += w + gap; placed += 1

            output_filename = unique_name("Passport_Photo", "jpg")
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            sheet.save(output_path, "JPEG", quality=95, dpi=(dpi, dpi))
            session["passport_photo_file"] = output_filename
            preview = sheet.copy()
            preview.thumbnail((800, 1000), PILImage.LANCZOS)
            buf = BytesIO_io.BytesIO()
            preview.save(buf, format="JPEG", quality=75)
            img_b64 = base64.b64encode(buf.getvalue()).decode()
            return render_template("passport_photo.html",
                                   preview_image=img_b64, download_ready=True,
                                   copies=placed, size=size_preset)
        except Exception as e:
            return render_template("passport_photo.html", error=f"Error: {str(e)}")
    return render_template("passport_photo.html")


@app.route("/passport-photo/download")
def passport_photo_download():
    output_filename = session.get("passport_photo_file")
    if not output_filename:
        return redirect("/passport-photo")
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)
    if os.path.exists(output_path):
        return send_file(output_path, as_attachment=True, download_name="Passport_Photo.jpg")
    return redirect("/passport-photo")


@app.route("/favicon-generator")
def favicon_generator_page():
    return render_template("favicon_generator.html")


@app.route("/api/favicon-generator", methods=["POST"])
def api_favicon_generator():
    try:
        from PIL import Image as PILImage
        import zipfile

        cleanup_old_favicons()

        if "image" not in request.files:
            return jsonify({"error": "No image uploaded"}), 400

        file = request.files["image"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400

        try:
            img = PILImage.open(file.stream).convert("RGBA")
        except Exception:
            return jsonify({"error": "Could not read image. Please upload a valid PNG/JPG/WEBP."}), 400

        FAVICON_SIZES = {
            "favicon-16x16.png": 16,
            "favicon-32x32.png": 32,
            "favicon-48x48.png": 48,
            "apple-touch-icon.png": 180,
            "android-chrome-192x192.png": 192,
            "android-chrome-512x512.png": 512,
        }

        # Make it square by padding onto a transparent canvas (no cropping)
        max_side = max(img.size)
        square_img = PILImage.new("RGBA", (max_side, max_side), (0, 0, 0, 0))
        square_img.paste(img, ((max_side - img.width) // 2, (max_side - img.height) // 2))

        batch_id = uuid.uuid4().hex[:10]
        batch_dir = os.path.join(FAVICON_OUTPUT_DIR, batch_id)
        os.makedirs(batch_dir, exist_ok=True)

        previews = []
        png_paths = {}

        for filename, size in FAVICON_SIZES.items():
            resized = square_img.resize((size, size), PILImage.LANCZOS)
            out_path = os.path.join(batch_dir, filename)
            resized.save(out_path, "PNG")
            png_paths[filename] = out_path
            previews.append({
                "url": url_for("static", filename=f"generated/favicons/{batch_id}/{filename}"),
                "label": filename.replace(".png", "")
            })

        ico_path = os.path.join(batch_dir, "favicon.ico")
        ico_source = square_img.resize((256, 256), PILImage.LANCZOS)
        ico_source.save(ico_path, format="ICO", sizes=[(16, 16), (32, 32), (48, 48)])

        manifest_content = """{
  "name": "",
  "short_name": "",
  "icons": [
    {
      "src": "/android-chrome-192x192.png",
      "sizes": "192x192",
      "type": "image/png"
    },
    {
      "src": "/android-chrome-512x512.png",
      "sizes": "512x512",
      "type": "image/png"
    }
  ],
  "theme_color": "#ffffff",
  "background_color": "#ffffff",
  "display": "standalone"
}"""
        manifest_path = os.path.join(batch_dir, "site.webmanifest")
        with open(manifest_path, "w") as f:
            f.write(manifest_content)

        zip_filename = f"favicon-package-{batch_id}.zip"
        zip_path = os.path.join(batch_dir, zip_filename)
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.write(ico_path, "favicon.ico")
            for filename, path in png_paths.items():
                zf.write(path, filename)
            zf.write(manifest_path, "site.webmanifest")

        html_snippet = """<link rel="icon" type="image/x-icon" href="/favicon.ico">
<link rel="icon" type="image/png" sizes="16x16" href="/favicon-16x16.png">
<link rel="icon" type="image/png" sizes="32x32" href="/favicon-32x32.png">
<link rel="apple-touch-icon" sizes="180x180" href="/apple-touch-icon.png">
<link rel="icon" type="image/png" sizes="192x192" href="/android-chrome-192x192.png">
<link rel="manifest" href="/site.webmanifest">"""

        return jsonify({
            "previews": previews,
            "zip_url": url_for("static", filename=f"generated/favicons/{batch_id}/{zip_filename}"),
            "html_snippet": html_snippet
        })
    except Exception as e:
        return jsonify({"error": f"Error: {str(e)}"}), 500


@app.route("/image-protect", methods=["GET", "POST"])
def image_protect():
    """Watermark an image and lock it into a password-protected PDF."""
    if request.method == "POST":
        try:
            from PIL import Image as PILImage, ImageDraw, ImageFont

            img_file = request.files.get("image")
            watermark_text = request.form.get("watermark_text", "").strip()
            password = request.form.get("password", "").strip()

            if not img_file or img_file.filename == "":
                return render_template("image_protect.html", error="Please select an image.")
            if not password:
                return render_template("image_protect.html", error="Please set a password to protect the file.")

            img = PILImage.open(img_file).convert("RGB")

            # Apply a tiled diagonal watermark if the user provided text
            if watermark_text:
                img = img.convert("RGBA")
                watermark_layer = PILImage.new("RGBA", img.size, (255, 255, 255, 0))
                draw = ImageDraw.Draw(watermark_layer)

                font_size = max(20, img.width // 18)
                try:
                    font = ImageFont.truetype("arial.ttf", font_size)
                except Exception:
                    font = ImageFont.load_default()

                bbox = draw.textbbox((0, 0), watermark_text, font=font)
                text_w = bbox[2] - bbox[0]
                text_h = bbox[3] - bbox[1]
                step_x = text_w + 120
                step_y = text_h + 120

                for y in range(-img.height, img.height * 2, step_y):
                    for x in range(-img.width, img.width * 2, step_x):
                        draw.text((x, y), watermark_text, font=font, fill=(255, 255, 255, 110))

                watermark_layer = watermark_layer.rotate(30, expand=0)
                img = PILImage.alpha_composite(img, watermark_layer).convert("RGB")

            # Convert the (watermarked) image to a PDF
            unprotected_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_to_protect.pdf")
            img.save(unprotected_path, "PDF", resolution=150.0)

            # Password-protect the PDF
            reader = PdfReader(unprotected_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(password)

            output_filename = unique_name("Protected_Image", "pdf")
            output_path = os.path.join(OUTPUT_FOLDER, output_filename)
            with open(output_path, "wb") as f:
                writer.write(f)

            session["protected_image_file"] = output_filename
            return render_template("image_protect.html", success=True)
        except Exception as e:
            return render_template("image_protect.html", error=f"Error: {str(e)}")
    return render_template("image_protect.html")


@app.route("/image-protect/download")
def image_protect_download():
    output_filename = session.get("protected_image_file")
    if not output_filename:
        return redirect("/image-protect")
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)
    if os.path.exists(output_path):
        return send_file(output_path, as_attachment=True, download_name="Protected_Image.pdf")
    return redirect("/image-protect")


# ──────────────────────────────────────────
# CALCULATORS
# ──────────────────────────────────────────

@app.route("/age-calculator")
def age_calculator():
    return render_template("age.html")

@app.route("/calculators/age")
def age_calculator_old():
    return redirect("/age-calculator", code=301)

@app.route("/bmi-calculator")
def bmi_calculator():
    return render_template("bmi.html")

@app.route("/calculators/bmi")
def bmi_calculator_old():
    return redirect("/bmi-calculator", code=301)

@app.route("/gst-calculator")
def gst_calculator():
    return render_template("gst.html")

@app.route("/calculators/gst")
def gst_calculator_old():
    return redirect("/gst-calculator", code=301)

@app.route("/emi-calculator")
def emi_calculator():
    return render_template("emi.html")

@app.route("/calculators/emi")
def emi_calculator_old():
    return redirect("/emi-calculator", code=301)

@app.route("/percentage-calculator")
def percentage_calculator():
    return render_template("percentage.html")

@app.route("/calculators/percentage")
def percentage_calculator_old():
    return redirect("/percentage-calculator", code=301)

@app.route("/sip-calculator")
def sip_calculator():
    return render_template("sip_calculator.html")

@app.route("/calculators/sip")
def sip_calculator_old():
    return redirect("/sip-calculator", code=301)

@app.route("/date-calculator")
def date_calculator():
    return render_template("date_calculator.html")

@app.route("/calculators/date")
def date_calculator_old():
    return redirect("/date-calculator", code=301)

@app.route("/calorie-calculator")
def calorie_calculator():
    return render_template("calorie_calculator.html")

@app.route("/calculators/calorie")
def calorie_calculator_old():
    return redirect("/calorie-calculator", code=301)

@app.route("/tip-calculator")
def tip_calculator():
    return render_template("tip_calculator.html")

@app.route("/tools/tip-calculator")
def tip_calculator_old():
    return redirect("/tip-calculator", code=301)

@app.route("/amount-in-words")
def amount_in_words_page():
    return render_template("amount_in_words.html")


# ──────────────────────────────────────────
# NEW CALCULATORS — Land Area, Brick, Salary, HRA, Timestamp
# ──────────────────────────────────────────

@app.route("/land-area-converter")
def land_area_converter():
    return render_template("land_area_converter.html")

@app.route("/brick-calculator")
def brick_calculator():
    return render_template("brick_calculator.html")

@app.route("/in-hand-salary-calculator")
def in_hand_salary_calculator():
    return render_template("in_hand_salary_calculator.html")

@app.route("/hra-exemption-calculator")
def hra_exemption_calculator():
    return render_template("hra_exemption_calculator.html")

@app.route("/timestamp-converter")
def timestamp_converter():
    return render_template("timestamp_converter.html")


# ──────────────────────────────────────────
# TEXT TOOLS
# ──────────────────────────────────────────

@app.route("/word-counter")
def word_counter():
    return render_template("word_counter.html")

@app.route("/tools/word-counter")
def word_counter_old():
    return redirect("/word-counter", code=301)

@app.route("/case-converter")
def case_converter():
    return render_template("case_converter.html")

@app.route("/tools/case-converter")
def case_converter_old():
    return redirect("/case-converter", code=301)

@app.route("/mb-converter")
def mb_converter():
    return render_template("mb_converter.html")

@app.route("/tools/mb-converter")
def mb_converter_old():
    return redirect("/mb-converter", code=301)

@app.route("/lorem-ipsum")
def lorem_generator():
    return render_template("lorem.html")

@app.route("/tools/lorem")
def lorem_generator_old():
    return redirect("/lorem-ipsum", code=301)

@app.route("/qr-generator")
def qr_generator():
    return render_template("qr_generator.html")

@app.route("/tools/qr-generator")
def qr_generator_old():
    return redirect("/qr-generator", code=301)

@app.route("/password-generator")
def password_generator():
    return render_template("password_generator.html")

@app.route("/tools/password-generator")
def password_generator_old():
    return redirect("/password-generator", code=301)

@app.route("/currency-converter")
def currency_converter():
    return render_template("currency_converter.html")

@app.route("/tools/currency-converter")
def currency_converter_old():
    return redirect("/currency-converter", code=301)

@app.route("/unit-converter")
def unit_converter():
    return render_template("unit_converter.html")

@app.route("/tools/unit-converter")
def unit_converter_old():
    return redirect("/unit-converter", code=301)

@app.route("/cft-calculator")
def cft_calculator():
    return render_template("cft_calculator.html")

@app.route("/sand-calculator")
def sand_calculator():
    return render_template("sand_calculator.html")


# ──────────────────────────────────────────
# BUSINESS TOOLS
# ──────────────────────────────────────────

@app.route("/invoice-generator", methods=["GET", "POST"])
def invoice_generator():
    if request.method == "POST":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Spacer, Paragraph
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            from reportlab.lib.units import mm
            from datetime import datetime as dt

            from_name    = request.form.get("from_name", "")
            from_address = request.form.get("from_address", "")
            from_gst     = request.form.get("from_gst", "")
            to_name      = request.form.get("to_name", "")
            to_address   = request.form.get("to_address", "")
            invoice_no   = request.form.get("invoice_no", "INV-001")
            invoice_date = request.form.get("invoice_date", dt.now().strftime("%d/%m/%Y"))
            due_date     = request.form.get("due_date", "")
            gst_rate     = float(request.form.get("gst_rate", 18))
            notes        = request.form.get("notes", "Thank you for your business!")

            items_desc = request.form.getlist("item_desc[]")
            items_qty  = request.form.getlist("item_qty[]")
            items_rate = request.form.getlist("item_rate[]")

            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Invoice", "pdf"))
            doc = SimpleDocTemplate(output_path, pagesize=A4,
                                    leftMargin=15*mm, rightMargin=15*mm,
                                    topMargin=15*mm, bottomMargin=15*mm)
            styles = getSampleStyleSheet()
            blue  = colors.HexColor('#2563EB')
            light = colors.HexColor('#F9FAFB')
            elements = []

            header_data = [[
                Paragraph(f"<font size='20' color='#2563EB'><b>INVOICE</b></font>", styles['Normal']),
                Paragraph(f"<font size='9' color='#6B7280'>Invoice No: <b>{invoice_no}</b><br/>Date: {invoice_date}" +
                          (f"<br/>Due: {due_date}" if due_date else "") + "</font>", styles['Normal'])
            ]]
            header_table = Table(header_data, colWidths=[90*mm, 80*mm])
            header_table.setStyle(TableStyle([('VALIGN', (0,0), (-1,-1), 'TOP'), ('ALIGN', (1,0), (1,0), 'RIGHT')]))
            elements.append(header_table)
            elements.append(Spacer(1, 8*mm))

            ft_data = [[
                Paragraph(f"<font size='8' color='#6B7280'>FROM</font><br/><font size='10'><b>{from_name}</b></font><br/>"
                          f"<font size='8' color='#6B7280'>{from_address.replace(chr(10),'<br/>')}</font>" +
                          (f"<br/><font size='8' color='#6B7280'>GST: {from_gst}</font>" if from_gst else ""), styles['Normal']),
                Paragraph(f"<font size='8' color='#6B7280'>BILL TO</font><br/><font size='10'><b>{to_name}</b></font><br/>"
                          f"<font size='8' color='#6B7280'>{to_address.replace(chr(10),'<br/>')}</font>", styles['Normal'])
            ]]
            ft_table = Table(ft_data, colWidths=[85*mm, 85*mm])
            ft_table.setStyle(TableStyle([
                ('VALIGN', (0,0), (-1,-1), 'TOP'), ('BACKGROUND', (0,0), (-1,-1), light),
                ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor('#E5E7EB')),
                ('PADDING', (0,0), (-1,-1), 8), ('LINEAFTER', (0,0), (0,-1), 0.5, colors.HexColor('#E5E7EB')),
            ]))
            elements.append(ft_table)
            elements.append(Spacer(1, 8*mm))

            item_rows = [['#', 'Description', 'Qty', 'Rate (₹)', 'Amount (₹)']]
            subtotal = 0
            for i, (desc, qty, rate) in enumerate(zip(items_desc, items_qty, items_rate)):
                if desc.strip():
                    q = float(qty) if qty else 1
                    r = float(rate) if rate else 0
                    amount = q * r
                    subtotal += amount
                    item_rows.append([str(i+1), desc, f"{q:.0f}", f"₹{r:,.2f}", f"₹{amount:,.2f}"])
            gst_amount = subtotal * gst_rate / 100
            total = subtotal + gst_amount
            item_rows.append(['', '', '', 'Subtotal', f'₹{subtotal:,.2f}'])
            item_rows.append(['', '', '', f'GST ({gst_rate:.0f}%)', f'₹{gst_amount:,.2f}'])
            item_rows.append(['', '', '', Paragraph('<b>TOTAL</b>', styles['Normal']),
                              Paragraph(f'<b>₹{total:,.2f}</b>', styles['Normal'])])

            last = len(item_rows) - 1
            items_table = Table(item_rows, colWidths=[10*mm, 80*mm, 15*mm, 30*mm, 35*mm])
            items_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), blue), ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'), ('FONTSIZE', (0,0), (-1,-1), 8),
                ('ROWBACKGROUNDS', (0,1), (-1, last-3), [colors.white, light]),
                ('GRID', (0,0), (-1, last-3), 0.3, colors.HexColor('#E5E7EB')),
                ('ALIGN', (2,0), (-1,-1), 'RIGHT'), ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('PADDING', (0,0), (-1,-1), 5),
                ('LINEABOVE', (3, last-2), (-1, last), 0.5, colors.HexColor('#E5E7EB')),
                ('BACKGROUND', (3, last), (-1, last), colors.HexColor('#EFF6FF')),
                ('FONTNAME', (3, last), (-1, last), 'Helvetica-Bold'),
            ]))
            elements.append(items_table)
            elements.append(Spacer(1, 8*mm))
            if notes:
                elements.append(Paragraph(f"<font size='8' color='#6B7280'><b>Notes:</b> {notes}</font>", styles['Normal']))
            doc.build(elements)
            return send_file(output_path, as_attachment=True, download_name=f"Invoice_{invoice_no}.pdf")
        except Exception as e:
            return render_template("invoice_generator.html", error=f"Error: {str(e)}")
    return render_template("invoice_generator.html")


# ──────────────────────────────────────────
# EXTRA TOOLS
# ──────────────────────────────────────────
@app.route("/youtube-thumbnail")
def youtube_thumbnail():
    return render_template("youtube_thumbnail.html")

@app.route("/income-tax-calculator")
def income_tax_calculator():
    return render_template("income_tax_calculator.html")

@app.route("/json-formatter")
def json_formatter():
    return render_template("json_formatter.html")

@app.route("/color-picker")
def color_picker():
    return render_template("color_picker.html")

@app.route("/base64-url-encoder")
def base64_url_encoder_page():
    return render_template("base64_url_encoder.html")


# ──────────────────────────────────────────
# NEWEST TOOLS — Signature, Meme, Resume, Page Numbers
# ──────────────────────────────────────────

@app.route("/signature-generator")
def signature_generator():
    return render_template("signature_generator.html")


@app.route("/meme-generator")
def meme_generator():
    return render_template("meme_generator.html")


@app.route("/resume-builder", methods=["GET", "POST"])
def resume_builder():
    if request.method == "POST":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            from reportlab.lib.units import mm

            full_name = request.form.get("full_name", "")
            email     = request.form.get("email", "")
            phone     = request.form.get("phone", "")
            location  = request.form.get("location", "")
            link      = request.form.get("link", "")
            summary   = request.form.get("summary", "")
            skills    = request.form.get("skills", "")

            exp_titles    = request.form.getlist("exp_title[]")
            exp_companies = request.form.getlist("exp_company[]")
            exp_durations = request.form.getlist("exp_duration[]")
            exp_descs     = request.form.getlist("exp_desc[]")

            edu_degrees = request.form.getlist("edu_degree[]")
            edu_schools = request.form.getlist("edu_school[]")
            edu_years   = request.form.getlist("edu_year[]")

            if not full_name.strip():
                return render_template("resume_builder.html", error="Please enter your full name.")

            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Resume", "pdf"))
            doc = SimpleDocTemplate(output_path, pagesize=A4,
                                    leftMargin=18*mm, rightMargin=18*mm,
                                    topMargin=16*mm, bottomMargin=16*mm)
            styles = getSampleStyleSheet()
            elements = []

            elements.append(Paragraph(f"<font size='22'><b>{full_name}</b></font>", styles['Normal']))
            contact_bits = [x for x in [email, phone, location, link] if x.strip()]
            if contact_bits:
                elements.append(Spacer(1, 3*mm))
                elements.append(Paragraph(f"<font size='9' color='#6B7280'>{' | '.join(contact_bits)}</font>", styles['Normal']))
            elements.append(Spacer(1, 6*mm))

            if summary.strip():
                elements.append(Paragraph("<font size='11' color='#2563EB'><b>PROFESSIONAL SUMMARY</b></font>", styles['Normal']))
                elements.append(Spacer(1, 2*mm))
                elements.append(Paragraph(f"<font size='9.5'>{summary}</font>", styles['Normal']))
                elements.append(Spacer(1, 6*mm))

            if any(t.strip() for t in exp_titles):
                elements.append(Paragraph("<font size='11' color='#2563EB'><b>WORK EXPERIENCE</b></font>", styles['Normal']))
                elements.append(Spacer(1, 3*mm))
                for title, company, duration, desc in zip(exp_titles, exp_companies, exp_durations, exp_descs):
                    if not title.strip():
                        continue
                    elements.append(Paragraph(f"<font size='10'><b>{title}</b></font> — <font size='10' color='#6B7280'>{company}</font>", styles['Normal']))
                    if duration.strip():
                        elements.append(Paragraph(f"<font size='8.5' color='#6B7280'>{duration}</font>", styles['Normal']))
                    if desc.strip():
                        elements.append(Spacer(1, 1*mm))
                        elements.append(Paragraph(f"<font size='9'>{desc}</font>", styles['Normal']))
                    elements.append(Spacer(1, 4*mm))

            if any(d.strip() for d in edu_degrees):
                elements.append(Paragraph("<font size='11' color='#2563EB'><b>EDUCATION</b></font>", styles['Normal']))
                elements.append(Spacer(1, 3*mm))
                for degree, school, year in zip(edu_degrees, edu_schools, edu_years):
                    if not degree.strip():
                        continue
                    elements.append(Paragraph(f"<font size='10'><b>{degree}</b></font> — <font size='10' color='#6B7280'>{school}</font>", styles['Normal']))
                    if year.strip():
                        elements.append(Paragraph(f"<font size='8.5' color='#6B7280'>{year}</font>", styles['Normal']))
                    elements.append(Spacer(1, 4*mm))

            if skills.strip():
                elements.append(Paragraph("<font size='11' color='#2563EB'><b>SKILLS</b></font>", styles['Normal']))
                elements.append(Spacer(1, 2*mm))
                elements.append(Paragraph(f"<font size='9.5'>{skills}</font>", styles['Normal']))

            doc.build(elements)
            return send_file(output_path, as_attachment=True, download_name=f"{full_name.replace(' ', '_')}_Resume.pdf")
        except Exception as e:
            return render_template("resume_builder.html", error=f"Error: {str(e)}")
    return render_template("resume_builder.html")


@app.route("/add-page-numbers", methods=["GET", "POST"])
def add_page_numbers():
    if request.method == "POST":
        try:
            pdf = request.files.get("pdf")
            position = request.form.get("position", "bottom-center")
            start_number = int(request.form.get("start_number", 1))
            fmt = request.form.get("format", "number")

            if not pdf or pdf.filename == "":
                return render_template("add_page_numbers.html", error="Please select a PDF file.")

            input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4().hex[:8]}_{pdf.filename}")
            output_path = os.path.join(OUTPUT_FOLDER, unique_name("Numbered", "pdf"))
            pdf.save(input_path)

            doc = fitz.open(input_path)
            total = len(doc)

            for i, page in enumerate(doc):
                num = start_number + i
                if fmt == "page_of_total":
                    text = f"Page {num} of {total + start_number - 1}"
                elif fmt == "dash":
                    text = f"- {num} -"
                else:
                    text = str(num)

                rect = page.rect
                margin = 24
                font_size = 10

                if "bottom" in position:
                    y = rect.height - margin
                else:
                    y = margin + font_size

                if "center" in position:
                    x = rect.width / 2 - (len(text) * font_size / 4)
                elif "right" in position:
                    x = rect.width - margin - (len(text) * font_size / 2)
                else:
                    x = margin

                page.insert_text((x, y), text, fontsize=font_size, color=(0.2, 0.2, 0.2))

            doc.save(output_path)
            doc.close()
            return send_file(output_path, as_attachment=True, download_name="Numbered_" + pdf.filename)
        except Exception as e:
            return render_template("add_page_numbers.html", error=f"Error: {str(e)}")
    return render_template("add_page_numbers.html")


# ──────────────────────────────────────────
# BLOG
# ──────────────────────────────────────────

BLOG_POSTS = [
    {
        "slug": "merge-pdf-files-online",
        "title": "How to Merge PDF Files Online for Free (No Signup Required)",
        "excerpt": "Learn how to combine multiple PDF files into one document online for free, in 3 simple steps.",
        "category": "PDF Tools",
        "date": "July 1, 2026"
    },
    {
        "slug": "reduce-image-size-to-kb",
        "title": "How to Reduce Image Size to 50 KB or 100 KB Online (Free)",
        "excerpt": "Reduce any JPG or PNG photo to an exact KB size for exam forms and job applications — free and instant.",
        "category": "Image Tools",
        "date": "July 1, 2026"
    },
    {
        "slug": "compress-pdf-to-100kb-for-exam-form",
        "title": "How to Compress PDF to 100KB or 1MB Online for Free",
        "excerpt": "Reduce your PDF file size for SSC, Railway, Bank PO and other exam forms — free, instant, no signup required.",
        "category": "PDF Tools",
        "date": "July 3, 2026"
    },
    {
        "slug": "amount-in-words-for-cheque-gst-invoice",
        "title": "How to Write Amount in Words for Cheques & GST Invoices",
        "excerpt": "The correct format for cheque and invoice amounts in words, plus a free online converter for Indian Rupees.",
        "category": "Business Tools",
        "date": "July 3, 2026"
    },
    {
        "slug": "sign-pdf-online-free",
        "title": "How to Sign a PDF Online for Free (No Printing, No Signup)",
        "excerpt": "Sign a PDF online in seconds — draw or type your signature and place it anywhere on the document. Free, no printing, no signup.",
        "category": "PDF Tools",
        "date": "July 3, 2026"
    },
    {
    "slug": "how-to-make-a-meme-online-free",
    "title": "How to Make a Meme Online for Free — Fonts, Stickers & Filters",
    "excerpt": "Learn how to create a custom meme online for free — drag text anywhere, pick fonts, add stickers, apply filters, and download with no watermark.",
    "category": "Fun Tools",
    "date": "July 6, 2026"
    },
    {
    "slug": "exam-photo-signature-size-guide-2026",
    "title": "Government Exam Photo & Signature Size Guide 2026 (UPSC, SSC, IBPS, NEET, JEE)",
    "excerpt": "Interactive size checker plus exact KB and pixel specs for every major Indian government exam in 2026.",
    "category": "Image Tools",
    "date": "July 8, 2026"
    },
    {
    "slug": "split-pdf-crop-pdf-guide",
    "title": "Split PDF vs Crop PDF: Which One Do You Actually Need?",
    "excerpt": "Split PDF and Crop PDF solve different problems. Learn the difference, when to use each, and how to fix scanned PDFs, exam forms, and oversized files for free.",
    "category": "PDF Tools",
    "date": "July 12, 2026"
    },
    {
    "slug": "unlock-pdf-remove-password-online-free",
    "title": "Unlock PDF Online Free: How to Remove Password from a PDF",
    "excerpt": "Remove open and permission passwords from a PDF for free, no software. Learn the difference between password types and what to do if you've forgotten it.",
    "category": "PDF Tools",
    "date": "July 15, 2026"
    },
    {
    "slug": "itr-filing-2026-last-date-old-vs-new-tax-regime",
    "title": "ITR Filing 2026: Last Date, Which Form to File & Old vs New Tax Regime",
    "excerpt": "ITR filing last date is July 31, 2026 for salaried taxpayers. Which ITR form to choose, old vs new regime compared, and what happens if you miss the deadline.",
    "category": "Calculators",
    "date": "July 17, 2026"
    },
    {
    "slug": "home-loan-emi-calculation-explained",
    "title": "Home Loan EMI Calculation Explained (With a Real Example) — 2026 Guide",
    "excerpt": "The EMI formula explained with a real ₹40 lakh example, current 2026 interest rates, tax benefits under Section 24(b)/80C, and RBI's prepayment rules.",
    "category": "Calculators",
    "date": "July 19, 2026"
    },
    {
    "slug": "cement-sand-ratio-material-estimation-guide",
    "title": "Cement Sand Ratio Chart (1:4, 1:5, 1:6) + Material Calculator",
    "excerpt": "Standard cement-sand ratios for plastering, brickwork, and concrete, the dry volume factor everyone forgets, and a free calculator for exact bag-and-CFT quantities.",
    "category": "Calculators",
    "date": "July 19, 2026"
    }
    # Add more posts here as dicts, e.g.:
    # {
    #     "slug": "compress-pdf-guide",
    #     "title": "How to Compress a PDF Without Losing Quality",
    #     "excerpt": "...",
    #     "category": "PDF Tools",
    #     "date": "July 5, 2026"
    # },
]

@app.route("/blog")
def blog_index():
    return render_template("blog.html", posts=BLOG_POSTS)

@app.route("/blog/merge-pdf-files-online")
def blog_merge_pdf():
    return render_template("blog_post_merge_pdf.html")

@app.route("/blog/reduce-image-size-to-kb")
def blog_reduce_image_size():
    return render_template("blog_post_reduce_image_size.html")

@app.route("/blog/compress-pdf-to-100kb-for-exam-form")
def blog_compress_pdf_exam_form():
    return render_template("blog_post_compress_pdf_exam_form.html")

@app.route("/blog/amount-in-words-for-cheque-gst-invoice")
def blog_amount_in_words():
    return render_template("blog_post_amount_in_words.html")

@app.route("/blog/sign-pdf-online-free")
def blog_sign_pdf():
    return render_template("blog_post_sign_pdf.html")

@app.route("/blog/how-to-make-a-meme-online-free")
def blog_how_to_make_a_meme():
    return render_template("blog_post_how_to_make_a_meme.html")

@app.route("/blog/exam-photo-signature-size-guide-2026")
def blog_exam_photo_signature_size_guide():
    return render_template("blog_post_exam_photo_signature_size_guide.html")

@app.route("/blog/split-pdf-crop-pdf-guide")
def blog_split_crop_pdf_guide():
    return render_template("blog_post_split_crop_pdf_guide.html")

@app.route("/blog/unlock-pdf-remove-password-online-free")
def blog_unlock_pdf():
    return render_template("blog_post_unlock_pdf_guide.html")

@app.route("/blog/itr-filing-2026-last-date-old-vs-new-tax-regime")
def blog_itr_filing_2026():
    return render_template("blog_post_itr_filing_2026.html")

@app.route("/author/ekramul-hoque")
def author_ekramul_hoque():
    return render_template("author_ekramul_hoque.html", posts=BLOG_POSTS)

@app.route("/blog/home-loan-emi-calculation-explained")
def blog_home_loan_emi():
    return render_template("blog_post_home_loan_emi.html")

@app.route("/blog/cement-sand-ratio-material-estimation-guide")
def blog_cement_sand_ratio_guide():
    return render_template("blog_post_cement_sand_ratio_guide.html")

# NOTE: as you add more blog posts, add a matching @app.route("/blog/<slug>")
# function for each new HTML file, the same way as blog_merge_pdf() above.
# Also remember to add the new URL to the sitemap() function above, and
# add a LASTMOD_OVERRIDES entry with today's date for the new blog URL.


# ──────────────────────────────────────────
# RUN
# ──────────────────────────────────────────
if __name__ == "__main__":
    app.run(debug=True)